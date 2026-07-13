#!/usr/bin/env python3
"""進捗報告デッキ（pptx）-> 正規化markdown。

スライドタイトル=見出し、本文テキスト、スピーカーノート、表を抽出する。
画像・スクリーンショットは raw/ に原本が残るため内容は抽出せず、
枚数のみ記録する（Confluenceへは非公開 — docs/llm-wiki.md §7.2 / §9.3）。

判断（ステータス表・リスク一覧・マイルストーンの意味づけ等）は行わない。
構造の正規化のみ。テーブルは汎用的にmarkdown表として書き出す。
"""
from __future__ import annotations

import argparse
import sys
from datetime import datetime
from pathlib import Path


def load_presentation(path: Path):
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "python-pptx が必要です（pip install python-pptx）。"
            "requirements.txt を参照してください。"
        ) from e
    return Presentation(str(path))


def _picture_shape_type():
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    return MSO_SHAPE_TYPE.PICTURE


def shape_text_lines(shape) -> list[str]:
    lines = []
    if not shape.has_text_frame:
        return lines
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs) or paragraph.text
        text = text.strip()
        if text:
            lines.append(text)
    return lines


def table_to_markdown(shape) -> list[str]:
    table = shape.table
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return []
    lines = ["| " + " | ".join(rows[0]) + " |"]
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return lines


def extract_slide(slide, index: int, picture_type) -> list[str]:
    lines = []
    title = None
    if slide.shapes.title is not None:
        title = slide.shapes.title.text.strip() or None
    lines.append(f"## Slide {index}: {title or '(タイトルなし)'}")
    lines.append("")

    image_count = 0
    body_lines: list[str] = []
    table_blocks: list[list[str]] = []

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.shape_type == picture_type:
            image_count += 1
            continue
        if shape.has_table:
            table_blocks.append(table_to_markdown(shape))
            continue
        body_lines.extend(shape_text_lines(shape))

    if body_lines:
        lines.append("### 本文")
        lines.extend(f"- {line}" for line in body_lines)
        lines.append("")

    for i, block in enumerate(table_blocks, start=1):
        lines.append(f"### 表 {i}")
        lines.extend(block)
        lines.append("")

    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            lines.append("### スピーカーノート")
            lines.append(notes_text)
            lines.append("")

    if image_count:
        lines.append(f"（画像 {image_count} 件あり — Confluence非公開、raw/decksの原本を参照）")
        lines.append("")

    return lines


def render_markdown(pptx_path: Path, prs) -> str:
    lines = [
        "---",
        "source_type: deck",
        f"original_file: {pptx_path.as_posix()}",
        f"extracted_at: {datetime.now().astimezone().isoformat(timespec='seconds')}",
        f"slide_count: {len(prs.slides)}",
        "---",
        "",
        f"# 進捗デッキ: {pptx_path.stem}",
        "",
    ]
    picture_type = _picture_shape_type()
    for i, slide in enumerate(prs.slides, start=1):
        lines.extend(extract_slide(slide, i, picture_type))
    return "\n".join(lines).rstrip() + "\n"


def extract(path: Path) -> str:
    prs = load_presentation(path)
    return render_markdown(path, prs)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path, help="pptxファイル")
    parser.add_argument("-o", "--output", type=Path, default=None, help="出力先md（省略時は入力と同じディレクトリ・同名.md）")
    args = parser.parse_args()

    if not args.input.exists():
        print(f"入力ファイルが見つかりません: {args.input}", file=sys.stderr)
        return 1

    try:
        markdown = extract(args.input)
    except RuntimeError as e:
        print(f"抽出に失敗しました: {e}", file=sys.stderr)
        return 1

    output_path = args.output or args.input.with_suffix(".md")
    output_path.write_text(markdown, encoding="utf-8")
    print(f"書き出しました: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
